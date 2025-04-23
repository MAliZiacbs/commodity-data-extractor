# document_processor.py

import PyPDF2
from pptx import Presentation
import tempfile
import os
import io

class DocumentProcessor:
    """Process different document formats to extract text"""
    
    @staticmethod
    def extract_text_from_pdf(pdf_file):
        """Extract text content from PDF file"""
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        
        # Extract text from each page
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += f"\n--- Page {page_num + 1} ---\n"
            text += page.extract_text() + "\n"
        
        return text
    
    @staticmethod
    def extract_text_from_pptx(pptx_file):
        """Extract text content from PowerPoint file"""
        # Save to temporary file (required for python-pptx)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            tmp.write(pptx_file.getvalue())
            tmp_path = tmp.name
        
        # Process the presentation
        presentation = Presentation(tmp_path)
        text = ""
        
        # Extract text from each slide
        for i, slide in enumerate(presentation.slides):
            text += f"\n--- Slide {i + 1} ---\n"
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text.strip() + "\n"
            
            # Handle tables specially
            for shape in slide.shapes:
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text.strip())
                        text += " | ".join(row_text) + "\n"
        
        # Clean up
        os.unlink(tmp_path)
        return text
    
    @staticmethod
    def process_document(file):
        """Process document based on file type"""
        filename = file.name.lower()
        
        if filename.endswith('.pdf'):
            return DocumentProcessor.extract_text_from_pdf(file), "PDF"
        elif filename.endswith(('.pptx', '.ppt')):
            return DocumentProcessor.extract_text_from_pptx(file), "PowerPoint"
        else:
            raise ValueError(f"Unsupported file format: {filename}")